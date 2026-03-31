"""Motor desacoplado para generación de presentaciones PPTX."""

import argparse
import os
import sys
from copy import deepcopy
from pathlib import Path

from lxml import etree

import yaml
from pptx import Presentation
from pptx.util import Pt


class MappingError(Exception):
    """Error en el mapeo entre content y mapping."""


class TemplateError(Exception):
    """Error en la validación de la plantilla."""


class PPTXGenerator:
    """Genera presentaciones PPTX a partir de plantilla + mapping + contenido."""

    def __init__(self):
        self.content = None
        self.mapping = None
        self.prs = None
        self.content_dir = None

    def load_config(self, content_path: str):
        """Carga content.yml y el mapping.yml referenciado."""
        content_path = Path(content_path).resolve()
        self.content_dir = content_path.parent

        with open(content_path, encoding="utf-8") as f:
            self.content = yaml.safe_load(f)

        map_ref = self.content["config"]["map_ref"]
        map_path = (self.content_dir / map_ref).resolve()

        with open(map_path, encoding="utf-8") as f:
            self.mapping = yaml.safe_load(f)

        template_file = self.mapping["template_file"]
        template_path = (map_path.parent / template_file).resolve()

        if not template_path.exists():
            raise TemplateError(f"Plantilla no encontrada: {template_path}")

        self.prs = Presentation(str(template_path))
        self._clear_existing_slides()
        self._clean_layout_ghosts()

    def _clear_existing_slides(self):
        """Elimina todos los slides existentes de la plantilla base."""
        xml_slides = self.prs.slides._sldIdLst
        slides = list(xml_slides)
        for sldId in slides:
            r_id = sldId.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if r_id:
                self.prs.part.drop_rel(r_id)
            xml_slides.remove(sldId)

    def _clean_layout_ghosts(self):
        """Elimina shapes fantasma no-placeholder de los layouts."""
        for layout in self.prs.slide_layouts:
            spTree = layout.shapes._spTree
            ghosts = []
            for shape in layout.shapes:
                if not shape.is_placeholder and shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if "editar el estilo" in text or "modificar el estilo" in text:
                        ghosts.append(shape._element)
            for elem in ghosts:
                spTree.remove(elem)

    def validate_template(self):
        """Verifica que el mapping coincida con la plantilla."""
        num_layouts = len(self.prs.slide_layouts)
        layouts = self.mapping.get("layouts", {})

        for name, layout_def in layouts.items():
            idx = layout_def["index"]
            if idx >= num_layouts:
                raise TemplateError(
                    f"Layout '{name}' tiene index {idx}, pero la plantilla "
                    f"solo tiene {num_layouts} layouts (0-{num_layouts - 1})."
                )

            layout = self.prs.slide_layouts[idx]
            ph_map = layout_def.get("placeholders", {})
            available_phs = {ph.placeholder_format.idx for ph in layout.placeholders}

            for alias, ph_idx in ph_map.items():
                if ph_idx not in available_phs:
                    raise TemplateError(
                        f"Layout '{name}' (index {idx}): placeholder '{alias}' "
                        f"con idx={ph_idx} no existe. "
                        f"Disponibles: {sorted(available_phs)}"
                    )

    def render(self):
        """Genera los slides según el manifiesto de contenido."""
        layouts_map = self.mapping["layouts"]
        slides_def = self.content.get("slides", [])

        for i, slide_def in enumerate(slides_def, 1):
            layout_name = slide_def["layout"]

            if layout_name not in layouts_map:
                raise MappingError(
                    f"Slide {i}: layout '{layout_name}' no existe en el mapping. "
                    f"Disponibles: {list(layouts_map.keys())}"
                )

            layout_cfg = layouts_map[layout_name]
            layout_idx = layout_cfg["index"]
            ph_map = layout_cfg.get("placeholders", {})
            slide_layout = self.prs.slide_layouts[layout_idx]
            slide = self.prs.slides.add_slide(slide_layout)

            config = self.content.get("config", {})
            slide_style = slide_def.get("style", {})

            data = slide_def.get("data", {})
            for alias, value in data.items():
                if alias not in ph_map:
                    raise MappingError(
                        f"Slide {i} (layout '{layout_name}'): "
                        f"placeholder '{alias}' no existe en el mapping. "
                        f"Disponibles: {list(ph_map.keys())}"
                    )

                ph_idx = ph_map[alias]
                try:
                    placeholder = slide.placeholders[ph_idx]
                except KeyError:
                    raise TemplateError(
                        f"Slide {i} (layout '{layout_name}'): "
                        f"placeholder idx={ph_idx} no encontrado en el slide."
                    )

                # Prioridad: style del slide > config global > plantilla
                font_size = slide_style.get(f"font_size_{alias}") \
                    or config.get(f"font_size_{alias}")
                self._set_text(placeholder, str(value).rstrip(), font_size=font_size)

    def _set_text(self, placeholder, text: str, font_size: int = None):
        """Inyecta texto en un placeholder. font_size en puntos (ej. 14, 18)."""
        tf = placeholder.text_frame
        lines = text.split("\n")
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        # Guardar copia del primer párrafo como base de formato
        base_p = deepcopy(tf.paragraphs[0]._p)

        # Limpiar todos los párrafos existentes
        for p_elem in list(tf._txBody.findall("a:p", nsmap)):
            tf._txBody.remove(p_elem)

        # Crear un párrafo por línea, copiando formato base
        for line in lines:
            new_p = deepcopy(base_p)
            # Desactivar bullets automáticos heredados del master
            pPr = new_p.find("a:pPr", nsmap)
            if pPr is None:
                pPr = etree.SubElement(new_p, f"{{{nsmap['a']}}}pPr")
                new_p.insert(0, pPr)
            # Quitar nivel para evitar herencia de bullets
            if "lvl" in pPr.attrib:
                del pPr.attrib["lvl"]
            # Agregar buNone para desactivar bullets
            for old_bu in pPr.findall("a:buNone", nsmap):
                pPr.remove(old_bu)
            etree.SubElement(pPr, f"{{{nsmap['a']}}}buNone")

            tf._txBody.append(new_p)
            from pptx.text.text import _Paragraph

            para = _Paragraph(new_p, tf)
            if para.runs:
                para.runs[0].text = line
                if font_size:
                    para.runs[0].font.size = Pt(font_size)
            else:
                para.text = line
                if font_size and para.runs:
                    para.runs[0].font.size = Pt(font_size)

    def _extract_font_props(self, paragraph):
        """Extrae propiedades de fuente de un párrafo como referencia."""
        props = {}
        if paragraph.runs:
            font = paragraph.runs[0].font
            props["bold"] = font.bold
            props["italic"] = font.italic
            props["size"] = font.size
            props["name"] = font.name
            props["color"] = font.color.rgb if font.color and font.color.rgb else None
        return props

    def save(self):
        """Guarda la presentación generada."""
        output_file = self.content["config"]["output_file"]
        output_path = (self.content_dir / output_file).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        print(f"Presentación generada: {output_path}")


def inspect_template(template_path: str):
    """Muestra layouts y placeholders de una plantilla PPTX."""
    prs = Presentation(template_path)

    print(f"\nPlantilla: {template_path}")
    print(f"Dimensiones: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"Layouts disponibles: {len(prs.slide_layouts)}")
    print(f"Slides existentes: {len(prs.slides)}")
    print("=" * 70)

    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n  Layout {i}: \"{layout.name}\"")
        print(f"  {'-' * 50}")

        placeholders = list(layout.placeholders)
        if not placeholders:
            print("    (sin placeholders)")
            continue

        print(f"  {'Idx':<6} {'Nombre':<30} {'Tipo':<15}")
        print(f"  {'---':<6} {'-----':<30} {'----':<15}")
        for ph in placeholders:
            fmt = ph.placeholder_format
            print(f"  {fmt.idx:<6} {ph.name:<30} {str(fmt.type):<15}")

    print()


def main():
    parser = argparse.ArgumentParser(
        description="Motor de generación de presentaciones PPTX"
    )
    subparsers = parser.add_subparsers(dest="command", help="Comando a ejecutar")

    # Subcomando: run
    run_parser = subparsers.add_parser("run", help="Generar presentación")
    run_parser.add_argument("content", help="Ruta al archivo content.yml")

    # Subcomando: inspect
    inspect_parser = subparsers.add_parser(
        "inspect", help="Inspeccionar plantilla PPTX"
    )
    inspect_parser.add_argument("template", help="Ruta al archivo .pptx")

    args = parser.parse_args()

    if args.command == "run":
        gen = PPTXGenerator()
        gen.load_config(args.content)
        gen.validate_template()
        gen.render()
        gen.save()

    elif args.command == "inspect":
        inspect_template(args.template)

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
